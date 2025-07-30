import os
import re
import shutil
from tqdm import tqdm
import docx
import pptx
import openpyxl
import PyPDF2

# --- Configuration ---
# Note: The source directory should contain .txt, .docx, .pptx, and .xlsx files.
SOURCE_DIR = "/home/sam/sandbox/rag/rag-llamacpp/docs/cleaned"
DEST_DIR = "/home/sam/sandbox/rag/rag-llamacpp/docs/final_for_rag"

# --- Text Extraction Functions ---

def extract_text_from_docx(filepath):
    """Extracts text from a .docx file."""
    doc = docx.Document(filepath)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(filepath):
    """Extracts text from a .pptx file."""
    prs = pptx.Presentation(filepath)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)

def extract_text_from_pdf(filepath):
    """Extracts text from a .pdf file using PyPDF2."""
    try:
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
    except Exception as e:
        print(f"Could not process PDF {filepath} with PyPDF2, skipping. Error: {e}")
        return ""

def extract_text_from_xlsx(filepath):
    """Extracts text from an .xlsx file, sheet by sheet."""
    workbook = openpyxl.load_workbook(filepath)
    text_content = []
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        text_content.append(f"--- Sheet: {sheet_name} ---")
        for row in sheet.iter_rows():
            row_text = [str(cell.value) for cell in row if cell.value is not None]
            if row_text:
                text_content.append("\t".join(row_text))
    return "\n".join(text_content)


# --- Text Cleaning Function ---

def clean_document_text(text):
    """
    Applies a series of cleaning rules to the document text.
    """
    # Rule 1: Remove dotloop signature verification lines
    text = re.sub(r'dotloop signature verification:.*', '', text, flags=re.IGNORECASE)

    # Rule 2: Remove lines that are just short GUID-like strings or timestamps
    text = re.sub(r'^[A-Z0-9-]{10,}\s*$', '', text, flags=re.MULTILINE)
    text = re.sub(r'^\d{2}/\d{2}/\d{2}\s+\d{1,2}:\d{2}\s+(AM|PM)\s+EDT\s*$', '', text, flags=re.MULTILINE)

    # Rule 3: Replace multiple newlines with a single newline
    text = re.sub(r'\n\s*\n', '\n', text)

    # Rule 4: Remove leading/trailing whitespace from each line
    lines = [line.strip() for line in text.split('\n')]
    text = '\n'.join(lines)

    # Rule 5: Remove any lines that are shorter than 5 characters (often OCR noise or artifacts)
    lines = [line for line in text.split('\n') if len(line.strip()) > 4]
    text = '\n'.join(lines)

    return text

# --- Main Processing Logic ---

def main():
    """
    Main function to process all supported documents in the source directory.
    """
    if not os.path.exists(SOURCE_DIR):
        print(f"Error: Source directory not found at '{SOURCE_DIR}'")
        return

    os.makedirs(DEST_DIR, exist_ok=True)
    print(f"Source directory: '{SOURCE_DIR}'")
    print(f"Destination directory: '{DEST_DIR}'")

    supported_extensions = (".txt", ".docx", ".pptx", ".xlsx", ".pdf")
    source_files = [f for f in os.listdir(SOURCE_DIR) if f.endswith(supported_extensions)]

    if not source_files:
        print(f"No supported files ({', '.join(supported_extensions)}) found in '{SOURCE_DIR}'.")
        return

    print(f"Found {len(source_files)} documents to process.")

    for filename in tqdm(source_files, desc="Processing documents"):
        source_path = os.path.join(SOURCE_DIR, filename)
        # Output file will always be .txt
        dest_filename = f"{filename}.txt"
        dest_path = os.path.join(DEST_DIR, dest_filename)
        original_text = ''

        try:
            if filename.endswith('.txt'):
                with open(source_path, 'r', encoding='utf-8', errors='ignore') as f:
                    original_text = f.read()
            elif filename.endswith('.docx'):
                original_text = extract_text_from_docx(source_path)
            elif filename.endswith('.pptx'):
                original_text = extract_text_from_pptx(source_path)
            elif filename.endswith('.xlsx'):
                original_text = extract_text_from_xlsx(source_path)
            elif filename.endswith('.pdf'):
                original_text = extract_text_from_pdf(source_path)

            cleaned_text = clean_document_text(original_text)

            with open(dest_path, 'w', encoding='utf-8') as f:
                f.write(cleaned_text)

        except Exception as e:
            print(f"Error processing file {filename}: {e}")

    print("\nDocument processing complete.")
    print(f"Cleaned text files are located in: {DEST_DIR}")

if __name__ == "__main__":
    main()
